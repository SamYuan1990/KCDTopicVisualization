#!/usr/bin/env python3
"""
PPT/PPTX to Markdown Converter
遍历目录，将 PPTX 文件转换为 Markdown 格式，提取文本、表格和图片。
对于 PPT（旧格式）将给出提示并跳过。
"""

import os
import sys
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_shape(shape):
    """从形状中提取文本（支持段落）"""
    if not shape.has_text_frame:
        return ""
    text = ""
    for paragraph in shape.text_frame.paragraphs:
        para_text = paragraph.text.strip()
        if para_text:
            # 简单处理：段落间用空行分隔
            text += para_text + "\n\n"
    return text.strip()


def extract_table_from_shape(shape):
    """将 PPT 表格转换为 Markdown 表格字符串"""
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return None
    table = shape.table
    rows = []
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            # 获取单元格文本，去除多余换行
            cell_text = cell.text_frame.text.replace("\n", " ").strip()
            row_cells.append(cell_text)
        rows.append(row_cells)
    if not rows:
        return ""

    # 构建 Markdown 表格
    md_table = []
    # 表头行
    md_table.append("| " + " | ".join(rows[0]) + " |")
    # 分隔行
    md_table.append("|" + "|".join([" --- " for _ in rows[0]]) + "|")
    # 数据行
    for row in rows[1:]:
        md_table.append("| " + " | ".join(row) + " |")
    return "\n".join(md_table)


def extract_images_from_slide(slide, slide_idx, output_img_dir):
    """提取幻灯片中的图片并保存，返回 Markdown 图片引用列表"""
    img_refs = []
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            # 获取图片扩展名（如 'png', 'jpg'）
            ext = image.ext
            # 构建文件名：slide_索引_shape_索引.扩展名
            img_filename = f"slide_{slide_idx}_shape_{shape_idx}.{ext}"
            img_path = output_img_dir / img_filename
            # 保存图片
            with open(img_path, "wb") as f:
                f.write(image.blob)
            # 生成 Markdown 图片引用（相对路径）
            img_refs.append(f"![图片]({img_filename})")
    return img_refs


def pptx_to_markdown(pptx_path, output_dir=None):
    """
    将 PPTX 文件转换为 Markdown 文件
    :param pptx_path: PPTX 文件路径
    :param output_dir: 输出目录，默认与源文件同目录
    :return: 生成的 Markdown 文件路径，失败返回 None
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        print(f"文件不存在: {pptx_path}")
        return None

    # 确定输出目录
    if output_dir is None:
        output_dir = pptx_path.parent
    else:
        output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # 创建图片子目录（在输出目录下）
    images_dir = output_dir / "images"
    images_dir.mkdir(exist_ok=True)

    # 加载 PPTX
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"无法解析 {pptx_path}: {e}")
        return None

    md_lines = []
    for slide_idx, slide in enumerate(prs.slides):
        # 添加幻灯片分隔（可选）
        md_lines.append(f"## Slide {slide_idx + 1}\n")

        # 提取文本
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                md_lines.append(text)
                md_lines.append("")  # 空行分隔

        # 提取表格
        for shape in slide.shapes:
            table_md = extract_table_from_shape(shape)
            if table_md:
                md_lines.append(table_md)
                md_lines.append("")

        # 提取图片
        img_refs = extract_images_from_slide(slide, slide_idx, images_dir)
        if img_refs:
            md_lines.extend(img_refs)
            md_lines.append("")

    # 生成 Markdown 文件
    md_filename = pptx_path.stem + ".md"
    md_path = output_dir / md_filename
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))

    print(f"已转换: {pptx_path} -> {md_path}")
    return md_path


def convert_ppt_file(ppt_path, output_dir=None):
    """
    处理单个 PPT 或 PPTX 文件
    如果是 PPTX，调用转换函数；如果是 PPT，提示无法转换（可扩展）
    """
    ppt_path = Path(ppt_path)
    ext = ppt_path.suffix.lower()
    if ext == ".pptx":
        return pptx_to_markdown(ppt_path, output_dir)
    elif ext == ".ppt":
        print(f"跳过旧版 PPT 文件（不支持）: {ppt_path}")
        return None
    else:
        print(f"跳过非 PPT/PPTX 文件: {ppt_path}")
        return None


def main():
    parser = argparse.ArgumentParser(description="将指定目录下的 PPT/PPTX 文件转换为 Markdown")
    parser.add_argument("input_dir", help="要遍历的目录")
    parser.add_argument("--output", "-o", default=None, help="输出目录（默认与源文件同目录）")
    args = parser.parse_args()

    input_dir = Path(args.input_dir)
    if not input_dir.is_dir():
        print(f"错误: {input_dir} 不是一个有效的目录")
        sys.exit(1)

    output_root = Path(args.output) if args.output else None

    # 遍历目录
    for root, dirs, files in os.walk(input_dir):
        for file in files:
            if file.lower().endswith((".ppt", ".pptx")):
                src_path = Path(root) / file
                # 如果指定了输出根目录，则保持相对目录结构
                if output_root:
                    rel_path = src_path.relative_to(input_dir)
                    out_dir = output_root / rel_path.parent
                else:
                    out_dir = src_path.parent
                convert_ppt_file(src_path, out_dir)


if __name__ == "__main__":
    main()